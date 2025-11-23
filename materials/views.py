from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response
from .models import Material
from .serializers import MaterialSerializer
from .segmentation import segment_text
from .summarization import summarize_text
from .quiz import generate_quiz_from_summary
from .education import generate_educational_insights
import re
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
import pytesseract

# Make sure Tesseract is correctly pointed
pytesseract.pytesseract.tesseract_cmd = r"C:\Users\Jamie\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"

def clean_text(text):
    # Remove emojis
    text = re.sub(r'[\U0001F600-\U0001F64F\U0001F300-\U0001F5FF\U0001F680-\U0001F6FF\U0001F1E0-\U0001F1FF]', '', text)
    # Remove short lines (less than 5 words)
    lines = [line for line in text.split('\n') if len(line.split()) > 4]
    return '\n'.join(lines)

# --- SEGMENTATION ---
@api_view(["POST"])
@permission_classes([IsAuthenticated])
def segment_view(request):
    text = request.data.get("text", "")
    title = request.data.get("title", "")
    material_id = request.data.get("material_id")  # Optional: existing material

    if not text.strip():
        return Response({"error": "No text provided"}, status=400)

    segmented = segment_text(text)

    if material_id:
        # Update existing material
        try:
            material = Material.objects.get(id=material_id, user=request.user)
            material.raw_text = text
            material.segmented_data = segmented
            if title:
                material.title = title
            material.save()
        except Material.DoesNotExist:
            # If material_id is invalid, create a new one
            material = Material.objects.create(
                user=request.user,
                title=title or "Untitled",
                raw_text=text,
                segmented_data=segmented
            )
    else:
        # No material_id, create new
        material = Material.objects.create(
            user=request.user,
            title=title or "Untitled",
            raw_text=text,
            segmented_data=segmented
        )

    return Response(MaterialSerializer(material).data)

# --- SUMMARIZATION ---
@api_view(["POST"])
@permission_classes([IsAuthenticated])
def summarize_view(request):
    text = request.data.get("text", "")
    material_id = request.data.get("material_id")

    if not text.strip():
        return Response({"error": "No text provided"}, status=400)

    # Generate the summary
    summary_obj = summarize_text(text)
    summary_text = summary_obj.get("summary") if isinstance(summary_obj, dict) else str(summary_obj)

    if material_id:
        try:
            # Update existing material
            material = Material.objects.get(id=material_id, user=request.user)
            material.summary_data = {"summary": summary_text}
            material.save()
        except Material.DoesNotExist:
            # Invalid ID: create new material
            material = Material.objects.create(
                user=request.user,
                raw_text=text,
                summary_data={"summary": summary_text},
                title="Untitled"
            )
    else:
        # No material_id: create new material
        material = Material.objects.create(
            user=request.user,
            raw_text=text,
            summary_data={"summary": summary_text},
            title="Untitled"
        )

    # Return the summary and material ID to frontend
    return Response({
        "summary": summary_text,
        "id": material.id
    })


# --- GET ALL MATERIALS ---
@api_view(["GET"])
@permission_classes([IsAuthenticated])
def materials_list(request):
    materials = Material.objects.filter(user=request.user).order_by("-created_at")
    serializer = MaterialSerializer(materials, many=True)
    return Response(serializer.data)

# --- MATERIAL DETAIL (GET + DELETE) ---
@api_view(["GET", "DELETE"])
@permission_classes([IsAuthenticated])
def material_detail(request, material_id):
    try:
        material = Material.objects.get(id=material_id, user=request.user)
    except Material.DoesNotExist:
        return Response({"error": "Material not found"}, status=404)

    if request.method == "GET":
        serializer = MaterialSerializer(material)
        return Response(serializer.data)
    elif request.method == "DELETE":
        material.delete()
        return Response({"message": "Material deleted successfully."}, status=204)

# --- EDUCATIONAL INSIGHTS ---
@api_view(["GET"])
@permission_classes([IsAuthenticated])
def educational_insights(request, material_id):
    try:
        material = Material.objects.get(id=material_id, user=request.user)
    except Material.DoesNotExist:
        return Response({"error": "Material not found"}, status=404)

    if not material.segmented_data:
        return Response({"error": "Material has no segmentation yet"}, status=400)

    insights = generate_educational_insights(material.segmented_data)
    return Response(insights)

# --- FILE UPLOAD (Does NOT create Material) ---
@api_view(["POST"])
@permission_classes([IsAuthenticated])
def upload_file(request):
    if request.FILES.get('file'):
        file = request.FILES['file']
        ext = file.name.split('.')[-1].lower()
        text = ""

        if ext == "pdf":
            with pdfplumber.open(file) as pdf:
                text = "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
        elif ext == "docx":
            doc = docx.Document(file)
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext == "txt":
            text = file.read().decode('utf-8')
        elif ext == "pptx":
            prs = Presentation(file)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext in ["jpg", "png"]:
            img = Image.open(file).convert("RGB")
            text = pytesseract.image_to_string(img)
        else:
            return Response({"error": "Unsupported file type"}, status=400)

        # Clean text
        cleaned = clean_text(text)

        # DO NOT create Material here, just return cleaned text
        return Response({"cleaned_text": cleaned})

    return Response({"error": "No file uploaded"}, status=400)


# --- QUIZ GENERATION ---
@api_view(["POST"])
@permission_classes([IsAuthenticated])
def generate_quiz(request, material_id):
    try:
        material = Material.objects.get(id=material_id, user=request.user)
    except Material.DoesNotExist:
        return Response({"error": "Material not found"}, status=404)

    # Extract the summary string from the dict
    summary_text = material.summary_data.get("summary", "") if isinstance(material.summary_data, dict) else ""

    if not summary_text.strip():
        return Response(
            {"error": "No summary available to generate a quiz."},
            status=400
        )

    # Generate quiz from the summary
    quiz_result = generate_quiz_from_summary(summary_text, num_questions=5)

    if not quiz_result:
        return Response(
            {"error": "Quiz could not be generated from the summary."},
            status=500
        )

    # Save quiz in the material
    material.quiz_data = quiz_result
    material.save()

    return Response({"quiz": quiz_result})
